# -*- coding: utf-8 -*-
"""Praesentations-Generator der ai-firma (V3 — Dark Violet Edition).

Aufruf:  python vorlagen/praesentation/baue_praesentation.py <konfig.json> <ziel.pptx>
(aus der Repo-Wurzel). Schema der Konfig: siehe ANLEITUNG.md.

V3: Schwarz + kraeftiges Violett (kein Rot), Verlaeufe (gradient fill), dunkle Flaechen,
moderne duenne Titel-Typo (Segoe UI Light) + Bahnschrift-Labels. Bewegung wird
nachtraeglich per PowerPoint-COM gesetzt (bewegung_setzen.ps1).
"""
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

BREITE, HOEHE = Inches(13.333), Inches(7.5)
# Fonts (alle auf Windows vorinstalliert)
FTITLE = "Segoe UI Light"      # Folien-Titel, gross, elegant-duenn
FWORD  = "Segoe UI Semibold"   # Wortmarke, Preise, Zahlen
FLABEL = "Bahnschrift"         # Eyebrow-Labels, technisch-modern
FBODY  = "Segoe UI"            # Fliesstext
WEISS = RGBColor.from_string("FFFFFF")


def farbe(hexwert):
    return RGBColor.from_string(hexwert.lstrip("#"))


def lade_konfig(pfad):
    with open(pfad, encoding="utf-8") as f:
        k = json.load(f)
    std = {
        "grund": "#0A0710", "grund2": "#170A22", "flaeche": "#191026",
        "flaeche2": "#241634", "violett": "#7C3AED", "violett_tief": "#4C1D95",
        "violett_hell": "#B98CFF", "text": "#F3EFFA", "text_dim": "#9E93B5",
        "linie": "#2E2440", "gut": "#2DD4BF", "knapp": "#F5B33D", "leer": "#D946EF",
    }
    for name, wert in std.items():
        k["farben"].setdefault(name, wert)
    k["_f"] = {name: farbe(wert) for name, wert in k["farben"].items()}
    return k


# ---------- Bau-Helfer ----------

def _grad(shape, c1, c2, angle=90):
    shape.fill.gradient()
    try:
        shape.fill.gradient_angle = angle
    except Exception:
        pass
    stops = shape.fill.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = c1
    stops[1].position = 1.0
    stops[1].color.rgb = c2


def leere_folie(prs, k):
    folie = prs.slides.add_slide(prs.slide_layouts[6])
    folie.background.fill.solid()
    folie.background.fill.fore_color.rgb = k["_f"]["grund"]
    r = folie.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, BREITE, HOEHE)
    _grad(r, k["_f"]["grund"], k["_f"]["grund2"], 65)
    r.line.fill.background()
    r.shadow.inherit = False
    return folie


def absatz_fuellen(p, text, groesse, farbe_rgb, fett, font):
    for i, teil in enumerate(str(text).split("**")):
        if not teil:
            continue
        lauf = p.add_run()
        lauf.text = teil
        lauf.font.size = Pt(groesse)
        lauf.font.color.rgb = farbe_rgb
        lauf.font.bold = fett or (i % 2 == 1)
        lauf.font.name = font


def textbox(folie, x, y, b, h, text, groesse, farbe_rgb, *, fett=False,
            font=FBODY, zentriert=False, zeilen=None, anker=MSO_ANCHOR.TOP,
            abstand=6):
    box = folie.shapes.add_textbox(x, y, b, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anker
    eintraege = zeilen if zeilen is not None else [(text, groesse, farbe_rgb, fett, font)]
    for i, (t, g, c, f, fo) in enumerate(eintraege):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER if zentriert else PP_ALIGN.LEFT
        p.space_after = Pt(abstand)
        absatz_fuellen(p, t, g, c, f, fo)
    return box


def kachel(folie, x, y, b, h, fuellung=None, rand=None, radius=0.06,
           randstaerke=1.0, grad=None, gradwinkel=90):
    form = folie.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, b, h)
    form.adjustments[0] = radius
    if grad is not None:
        _grad(form, grad[0], grad[1], gradwinkel)
    else:
        form.fill.solid()
        form.fill.fore_color.rgb = fuellung
    if rand:
        form.line.color.rgb = rand
        form.line.width = Pt(randstaerke)
    else:
        form.line.fill.background()
    form.shadow.inherit = False
    return form


def pfeil(folie, x, y, b, h, fuellung):
    form = folie.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, b, h)
    form.fill.solid()
    form.fill.fore_color.rgb = fuellung
    form.line.fill.background()
    form.shadow.inherit = False
    return form


def ring(folie, cx, cy, d, farbe_rgb, staerke=1.5):
    e = folie.shapes.add_shape(MSO_SHAPE.OVAL, cx - d // 2, cy - d // 2, d, d)
    e.fill.background()
    e.line.color.rgb = farbe_rgb
    e.line.width = Pt(staerke)
    e.shadow.inherit = False
    return e


def wortmarke(folie, k, x, y, groesse=18):
    box = folie.shapes.add_textbox(x, y, Inches(2.4), Inches(0.5))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r1 = p.add_run(); r1.text = k["kunde"].upper()
    r1.font.name = FWORD; r1.font.bold = True; r1.font.size = Pt(groesse)
    r1.font.color.rgb = k["_f"]["text"]
    r2 = p.add_run(); r2.text = " ●"
    r2.font.size = Pt(int(groesse * 0.55)); r2.font.color.rgb = k["_f"]["violett_hell"]
    return box


def kopfzeile(folie, k, titel, untertitel=None):
    kachel(folie, 0, 0, BREITE, Inches(0.1), None, radius=0.0,
           grad=(k["_f"]["violett"], k["_f"]["violett_tief"]), gradwinkel=0)
    textbox(folie, Inches(0.6), Inches(0.42), Inches(9.6), Inches(0.85),
            titel, 31, k["_f"]["text"], font=FTITLE)
    if untertitel:
        textbox(folie, Inches(0.6), Inches(1.12), Inches(9.6), Inches(0.5),
                untertitel, 13.5, k["_f"]["text_dim"])
    wortmarke(folie, k, Inches(10.35), Inches(0.45))


def eyebrow(folie, k, x, y, text, breite=Inches(3.2)):
    return textbox(folie, x, y, breite, Inches(0.4), text.upper(), 12.5,
                   k["_f"]["violett_hell"], fett=True, font=FLABEL)


# ---------- Gezeichnete Handy-Mockups (Dark UI) ----------

def _mini(folie, x, y, b, h, text, g, c, *, fett=False, zentriert=True,
          font=FBODY, anker=MSO_ANCHOR.MIDDLE):
    box = textbox(folie, x, y, b, h, text, g, c, fett=fett, font=font,
                  zentriert=zentriert, anker=anker, abstand=1)
    box.text_frame.margin_left = box.text_frame.margin_right = Inches(0.04)
    box.text_frame.margin_top = box.text_frame.margin_bottom = Inches(0.02)
    return box


def phone(folie, k, x, y, typ, b=Inches(2.9), h=Inches(5.0)):
    f = k["_f"]
    kachel(folie, x, y, b, h, f["grund"], rand=f["violett"], radius=0.11, randstaerke=1.5)
    sx, sy = x + Inches(0.12), y + Inches(0.12)
    sb, sh = b - Inches(0.24), h - Inches(0.24)
    kachel(folie, sx, sy, sb, sh, f["flaeche"], radius=0.09)
    _mini(folie, sx + Inches(0.14), sy + Inches(0.08), sb - Inches(0.28), Inches(0.3),
          k["kunde"].upper() + " ●", 10, f["text"], fett=True, zentriert=False, font=FWORD)
    cy = sy + Inches(0.48)
    cx = sx + Inches(0.14)
    cb = sb - Inches(0.28)
    if typ == "uebersicht":
        dr = (cb - Inches(0.12)) / 3
        for i, (n, l, c) in enumerate([("3", "gut", f["gut"]), ("3", "knapp", f["knapp"]), ("2", "leer", f["leer"])]):
            tx = cx + i * (dr + Inches(0.06))
            kachel(folie, tx, cy, dr, Inches(0.74), f["flaeche2"], rand=f["linie"])
            kachel(folie, tx, cy, dr, Inches(0.06), c, radius=0.0)
            _mini(folie, tx, cy + Inches(0.08), dr, Inches(0.4), n, 17, c, fett=True, font=FWORD)
            _mini(folie, tx, cy + Inches(0.46), dr, Inches(0.24), l, 7.5, f["text_dim"])
        cy += Inches(0.9)
        kachel(folie, cx, cy, cb, Inches(0.85), f["flaeche2"], rand=f["linie"])
        _mini(folie, cx + Inches(0.1), cy + Inches(0.06), cb - Inches(0.2), Inches(0.3),
              "Heute", 9, f["text"], fett=True, zentriert=False)
        _mini(folie, cx + Inches(0.1), cy + Inches(0.34), cb - Inches(0.2), Inches(0.45),
              "**14** Verkäufe   ·   **3** Bestellungen", 9, f["text"], zentriert=False)
        cy += Inches(1.0)
        kachel(folie, cx, cy, cb, Inches(1.5), f["flaeche2"], rand=f["violett"], randstaerke=1.5)
        _mini(folie, cx + Inches(0.1), cy + Inches(0.06), cb - Inches(0.2), Inches(0.3),
              "KI-Vorschlag", 9, f["violett_hell"], fett=True, zentriert=False)
        _mini(folie, cx + Inches(0.1), cy + Inches(0.36), cb - Inches(0.2), Inches(0.6),
              "Hoodie M in Zürich leer. 12 nachbestellen?", 8.5, f["text"], zentriert=False,
              anker=MSO_ANCHOR.TOP)
        bb2 = (cb - Inches(0.3)) / 2
        kachel(folie, cx + Inches(0.1), cy + Inches(1.02), bb2, Inches(0.36), None,
               grad=(f["violett"], f["violett_tief"]), gradwinkel=0, radius=0.28)
        _mini(folie, cx + Inches(0.1), cy + Inches(1.02), bb2, Inches(0.36), "Übernehmen", 8, WEISS, fett=True)
        kachel(folie, cx + Inches(0.2) + bb2, cy + Inches(1.02), bb2, Inches(0.36), f["flaeche"], rand=f["linie"], radius=0.28)
        _mini(folie, cx + Inches(0.2) + bb2, cy + Inches(1.02), bb2, Inches(0.36), "Ablehnen", 8, f["text_dim"])
    elif typ == "scan":
        kachel(folie, cx, cy, cb, Inches(0.5), f["flaeche2"], rand=f["linie"])
        kachel(folie, cx + Inches(0.1), cy + Inches(0.12), Inches(0.5), Inches(0.26), f["violett"], radius=0.5)
        _mini(folie, cx + Inches(0.68), cy, cb - Inches(0.7), Inches(0.5),
              "Test: WLAN ausgeschaltet", 8.5, f["text"], zentriert=False)
        cy += Inches(0.66)
        kachel(folie, cx, cy, cb, Inches(0.78), None, grad=(f["violett"], f["violett_tief"]),
               gradwinkel=90, radius=0.14)
        _mini(folie, cx, cy, cb, Inches(0.78), "Artikel scannen", 12, WEISS, fett=True)
        cy += Inches(0.94)
        kachel(folie, cx, cy, cb, Inches(0.6), f["flaeche2"], rand=f["gut"])
        _mini(folie, cx + Inches(0.1), cy, cb - Inches(0.2), Inches(0.6),
              "**Eingebucht:** 1x Tee Box-Logo M", 8.5, f["gut"], zentriert=False)
        cy += Inches(0.76)
        kachel(folie, cx, cy, cb, Inches(1.05), f["flaeche2"], rand=f["knapp"])
        _mini(folie, cx + Inches(0.1), cy + Inches(0.08), cb - Inches(0.2), Inches(0.9),
              "**Warteschlange: 1 Buchung.** Wird automatisch nachgebucht, sobald WLAN da ist.",
              8.5, f["knapp"], zentriert=False, anker=MSO_ANCHOR.TOP)
    elif typ == "chat":
        dialog = [("Wo liegt Hoodie M?", True), ("Zürich 0, Embrach 4. Regal B, Fach 3.", False),
                  ("Was ist knapp?", True), ("Cap Sport und Beanie Logo.", False)]
        for text, ich in dialog:
            bw = int(cb * 0.82)
            bx = cx + (cb - bw) if ich else cx
            hh = Inches(0.55)
            if ich:
                kachel(folie, bx, cy, bw, hh, None, grad=(f["violett"], f["violett_tief"]),
                       gradwinkel=0, radius=0.32)
            else:
                kachel(folie, bx, cy, bw, hh, f["flaeche2"], rand=f["linie"], radius=0.32)
            _mini(folie, bx + Inches(0.1), cy, bw - Inches(0.2), hh, text, 8.5,
                  WEISS if ich else f["text"], zentriert=False)
            cy += hh + Inches(0.14)
        kachel(folie, cx, sy + sh - Inches(0.55), cb, Inches(0.4), f["flaeche2"], rand=f["linie"], radius=0.3)
        _mini(folie, cx + Inches(0.12), sy + sh - Inches(0.55), cb - Inches(0.24), Inches(0.4),
              "Frage eingeben", 8.5, f["text_dim"], zentriert=False)


# ---------- Folien ----------

def slide_titel(prs, k):
    f = leere_folie(prs, k)
    kachel(f, Inches(-1), Inches(2.05), Inches(15), Inches(2.9), None,
           grad=(k["_f"]["grund"], k["_f"]["violett_tief"]), gradwinkel=25, radius=0.0)
    kachel(f, 0, HOEHE - Inches(0.1), BREITE, Inches(0.1), None, radius=0.0,
           grad=(k["_f"]["violett_tief"], k["_f"]["violett"]), gradwinkel=0)
    box = f.shapes.add_textbox(Inches(1.2), Inches(2.35), Inches(11), Inches(1.5))
    p = box.text_frame.paragraphs[0]
    r1 = p.add_run(); r1.text = k["kunde"].upper()
    r1.font.name = FWORD; r1.font.bold = True; r1.font.size = Pt(58); r1.font.color.rgb = k["_f"]["text"]
    r2 = p.add_run(); r2.text = " ●"
    r2.font.size = Pt(32); r2.font.color.rgb = k["_f"]["violett_hell"]
    kachel(f, Inches(1.25), Inches(3.75), Inches(1.4), Inches(0.05), k["_f"]["violett_hell"], radius=0.0)
    textbox(f, Inches(1.2), Inches(3.95), Inches(10.5), Inches(0.9),
            k["claim"], 23, k["_f"]["violett_hell"], font=FTITLE)
    textbox(f, Inches(1.2), Inches(6.4), Inches(10), Inches(0.6),
            k["untertitel_titelfolie"] + "   ·   " + k["datum"], 13, k["_f"]["text_dim"], font=FLABEL)


def slide_momente(prs, k):
    f = leere_folie(prs, k)
    m = k["momente"]
    kopfzeile(f, k, m["titel"], m.get("untertitel"))
    b = Inches(3.9); x = Inches(0.6)
    for eintrag in m["punkte"]:
        kachel(f, x, Inches(1.8), b, Inches(2.5), k["_f"]["flaeche"], rand=k["_f"]["linie"])
        kachel(f, x, Inches(1.8), Inches(0.08), Inches(2.5), k["_f"]["violett"], radius=0.0)
        textbox(f, x + Inches(0.3), Inches(2.05), b - Inches(0.58), Inches(2.1), "", 14,
                k["_f"]["text"],
                zeilen=[(eintrag["kopf"], 15, k["_f"]["violett_hell"], True, FWORD),
                        (eintrag["text"], 13, k["_f"]["text"], False, FBODY)])
        x += b + Inches(0.31)
    s = m.get("stat")
    if s:
        kachel(f, Inches(0.6), Inches(4.7), Inches(12.13), Inches(2.1), None,
               grad=(k["_f"]["flaeche"], k["_f"]["violett_tief"]), gradwinkel=30)
        textbox(f, Inches(1.0), Inches(5.15), Inches(3.9), Inches(1.2),
                s["wert"], 38, WEISS, fett=True, font=FWORD)
        textbox(f, Inches(4.8), Inches(5.1), Inches(7.6), Inches(0.6),
                s["text"], 16, WEISS, fett=True)
        bx, by, bb = Inches(4.8), Inches(5.85), Inches(6.6)
        kachel(f, bx, by, bb, Inches(0.3), k["_f"]["grund"], radius=0.5)
        kachel(f, bx, by, int(bb * s["anteil"]), Inches(0.3), k["_f"]["violett_hell"], radius=0.5)
        textbox(f, Inches(4.8), Inches(6.3), Inches(7.6), Inches(0.4),
                "Quelle: " + s["quelle"], 10.5, k["_f"]["text_dim"])
    elif m.get("band"):
        band = m["band"]
        kachel(f, Inches(0.6), Inches(4.7), Inches(12.13), Inches(2.1), None,
               grad=(k["_f"]["flaeche"], k["_f"]["violett_tief"]), gradwinkel=30)
        eyebrow(f, k, Inches(1.0), Inches(4.95), band.get("titel", ""), Inches(3.0))
        bb = Inches(3.75); x = Inches(1.0)
        for punkt in band["punkte"]:
            kachel(f, x, Inches(5.5), Inches(0.09), Inches(0.75), k["_f"]["violett_hell"], radius=0.5)
            textbox(f, x + Inches(0.3), Inches(5.5), bb - Inches(0.3), Inches(1.1),
                    punkt, 16, WEISS, fett=True)
            x += bb + Inches(0.15)


def slide_prinzipien(prs, k):
    f = leere_folie(prs, k)
    p = k["prinzipien"]
    kopfzeile(f, k, p["titel"], p.get("untertitel"))
    b = Inches(3.9); x = Inches(0.6)
    for i, eintrag in enumerate(p["punkte"], 1):
        kachel(f, x, Inches(2.0), b, Inches(3.9), k["_f"]["flaeche"], rand=k["_f"]["linie"])
        kachel(f, x + Inches(0.28), Inches(2.3), Inches(0.6), Inches(0.6), None,
               grad=(k["_f"]["violett"], k["_f"]["violett_tief"]), gradwinkel=45, radius=0.22)
        textbox(f, x + Inches(0.28), Inches(2.36), Inches(0.6), Inches(0.5),
                str(i), 20, WEISS, fett=True, zentriert=True, font=FWORD)
        textbox(f, x + Inches(0.28), Inches(3.15), b - Inches(0.56), Inches(2.6), "", 14,
                k["_f"]["text"],
                zeilen=[(eintrag["kopf"], 18, k["_f"]["text"], True, FTITLE),
                        (eintrag["text"], 13, k["_f"]["text_dim"], False, FBODY)])
        x += b + Inches(0.31)


def slide_business(prs, k):
    f = leere_folie(prs, k)
    a = k["abteilungen"]
    kopfzeile(f, k, a["titel"], a.get("untertitel"))
    eyebrow(f, k, Inches(0.6), Inches(1.72), a["links_titel"], Inches(5.6))
    bb, bh = Inches(2.6), Inches(0.85)
    for i, name in enumerate(a["liste"]):
        x = Inches(0.6) + (i % 2) * (bb + Inches(0.2))
        y = Inches(2.3) + (i // 2) * (bh + Inches(0.18))
        kachel(f, x, y, bb, bh, k["_f"]["flaeche"], rand=k["_f"]["linie"])
        textbox(f, x, y, bb, bh, name, 12.5, k["_f"]["text_dim"], fett=True,
                zentriert=True, anker=MSO_ANCHOR.MIDDLE)
    pfeil(f, Inches(6.35), Inches(4.0), Inches(1.0), Inches(0.7), k["_f"]["violett_hell"])
    kachel(f, Inches(7.7), Inches(2.3), Inches(5.03), Inches(3.98), None,
           grad=(k["_f"]["violett"], k["_f"]["violett_tief"]), gradwinkel=50)
    textbox(f, Inches(8.1), Inches(2.85), Inches(4.2), Inches(3.2), "", 16, WEISS,
            zeilen=[(a["rechts_titel"], 25, WEISS, True, FTITLE),
                    ("", 8, WEISS, False, FBODY),
                    (a["rechts_text"], 14.5, farbe("EADFFB"), False, FBODY)])
    textbox(f, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5),
            a.get("fusszeile", ""), 12, k["_f"]["text_dim"])


def slide_os_flow(prs, k):
    f = leere_folie(prs, k)
    o = k["os_flow"]
    kopfzeile(f, k, o["titel"], o.get("untertitel"))
    eyebrow(f, k, Inches(0.6), Inches(1.9), o["links_titel"], Inches(3.2))
    y = Inches(2.4)
    for name in o["systeme"]:
        kachel(f, Inches(0.6), y, Inches(3.0), Inches(0.75), k["_f"]["flaeche"], rand=k["_f"]["linie"])
        textbox(f, Inches(0.6), y, Inches(3.0), Inches(0.75), name, 13, k["_f"]["text"],
                fett=True, zentriert=True, anker=MSO_ANCHOR.MIDDLE)
        y += Inches(0.95)
    pfeil(f, Inches(3.85), Inches(3.72), Inches(0.85), Inches(0.5), k["_f"]["violett_hell"])
    kachel(f, Inches(4.95), Inches(2.7), Inches(3.5), Inches(2.6), None,
           grad=(k["_f"]["violett"], k["_f"]["violett_tief"]), gradwinkel=55, radius=0.12)
    ring(f, Inches(6.7), Inches(4.0), Inches(2.15), farbe("C9A9FF"), 1.25)
    textbox(f, Inches(5.1), Inches(3.35), Inches(3.2), Inches(1.6), "", 14, WEISS,
            zentriert=True,
            zeilen=[(o["os_name"], 17, WEISS, True, FWORD),
                    (o["os_sub"], 13, farbe("EADFFB"), False, FBODY)])
    pfeil(f, Inches(8.7), Inches(3.72), Inches(0.85), Inches(0.5), k["_f"]["violett_hell"])
    eyebrow(f, k, Inches(9.8), Inches(1.9), o["rechts_titel"], Inches(3.0))
    y = Inches(2.4)
    for eintrag in o["menschen"]:
        kachel(f, Inches(9.8), y, Inches(3.0), Inches(1.3), k["_f"]["flaeche"], rand=k["_f"]["violett"], randstaerke=1.25)
        textbox(f, Inches(9.95), y + Inches(0.18), Inches(2.7), Inches(1.0), "", 12, WEISS,
                zeilen=[(eintrag["wer"], 14, k["_f"]["violett_hell"], True, FWORD),
                        (eintrag["was"], 12, k["_f"]["text_dim"], False, FBODY)])
        y += Inches(1.55)
    textbox(f, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5),
            o.get("fusszeile", ""), 12.5, k["_f"]["text"], fett=True)


def slide_mockup(prs, k, m):
    f = leere_folie(prs, k)
    kopfzeile(f, k, m["titel"], m.get("untertitel"))
    phone(f, k, Inches(1.3), Inches(1.9), m["typ"])
    y = Inches(2.4)
    for punkt in m["punkte"]:
        kachel(f, Inches(5.2), y, Inches(0.12), Inches(0.95), None,
               grad=(k["_f"]["violett"], k["_f"]["violett_tief"]), gradwinkel=90, radius=0.5)
        textbox(f, Inches(5.6), y, Inches(7.1), Inches(1.1), punkt, 16, k["_f"]["text"])
        y += Inches(1.35)


def slide_schritte(prs, k):
    f = leere_folie(prs, k)
    s = k["schritte"]
    kopfzeile(f, k, s["titel"], s.get("untertitel"))
    n = len(s["punkte"])
    bb = Inches(3.05); x = Inches(0.6)
    for i, eintrag in enumerate(s["punkte"], 1):
        chv = f.shapes.add_shape(MSO_SHAPE.CHEVRON, x, Inches(2.1), bb, Inches(0.85))
        chv.adjustments[0] = 0.35
        if i == n:
            _grad(chv, k["_f"]["violett"], k["_f"]["violett_tief"], 0)
            chv.line.fill.background()
        else:
            chv.fill.solid(); chv.fill.fore_color.rgb = k["_f"]["flaeche"]
            chv.line.color.rgb = k["_f"]["linie"]; chv.line.width = Pt(1)
        chv.shadow.inherit = False
        textbox(f, x + Inches(0.35), Inches(2.22), bb - Inches(0.6), Inches(0.6),
                str(i) + "   " + eintrag["kopf"], 15, WEISS if i == n else k["_f"]["text"], fett=True, font=FWORD)
        textbox(f, x + Inches(0.1), Inches(3.25), bb - Inches(0.35), Inches(2.6),
                eintrag["text"], 13, k["_f"]["text_dim"])
        x += bb + Inches(0.03)
    textbox(f, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.6),
            s.get("fusszeile", ""), 13.5, k["_f"]["violett_hell"], fett=True)


def slide_modelle_visual(prs, k):
    f = leere_folie(prs, k)
    m = k["modelle"]
    kopfzeile(f, k, m["titel"], m.get("untertitel"))
    basis_y = Inches(6.45)
    hoehen = [Inches(2.7), Inches(3.2), Inches(3.7)]
    bb = Inches(3.7); x = Inches(0.75)
    for modell in m["liste"]:
        stil = modell["stil"]
        h = hoehen[["hell", "akzent", "dunkel"].index(stil)]
        y = basis_y - h
        if stil == "hell":
            kachel(f, x, y, bb, h, k["_f"]["flaeche"], rand=k["_f"]["linie"], radius=0.04)
            txt, sub, nm = k["_f"]["text"], k["_f"]["text_dim"], k["_f"]["violett_hell"]
        elif stil == "akzent":
            kachel(f, x, y, bb, h, None, grad=(k["_f"]["violett"], k["_f"]["violett_tief"]),
                   gradwinkel=60, radius=0.04)
            txt, sub, nm = WEISS, farbe("EADFFB"), WEISS
        else:
            kachel(f, x, y, bb, h, None, grad=(k["_f"]["flaeche2"], k["_f"]["grund"]),
                   gradwinkel=60, radius=0.04)
            kachel(f, x, y, bb, Inches(0.09), k["_f"]["violett_hell"], radius=0.0)
            txt, sub, nm = k["_f"]["text"], k["_f"]["text_dim"], k["_f"]["violett_hell"]
        textbox(f, x + Inches(0.3), y + Inches(0.25), bb - Inches(0.6), h - Inches(0.4), "", 13,
                txt,
                zeilen=[(modell["stufe"], 12, sub, True, FLABEL),
                        (modell["name"], 22, nm, True, FTITLE),
                        (modell["preis"], 27, txt, True, FWORD),
                        (m["preis_einheit"], 12, sub, False, FBODY),
                        ("", 6, sub, False, FBODY),
                        (modell["kurz"], 12.5, sub, False, FBODY)])
        if modell.get("empfohlen"):
            kachel(f, x + Inches(0.3), y - Inches(0.28), Inches(1.9), Inches(0.48), None,
                   grad=(k["_f"]["violett_hell"], k["_f"]["violett"]), gradwinkel=0)
            textbox(f, x + Inches(0.3), y - Inches(0.24), Inches(1.9), Inches(0.42),
                    m["empfehlung_etikett"], 11.5, WEISS, fett=True, zentriert=True, font=FWORD)
        x += bb + Inches(0.35)
    textbox(f, Inches(0.75), Inches(6.65), Inches(12), Inches(0.5),
            m.get("fusszeile", ""), 11.5, k["_f"]["text_dim"])


def slide_klartext(prs, k, modell, einheit):
    f = leere_folie(prs, k)
    kopfzeile(f, k, modell["name"],
              modell["preis"] + " " + einheit + ("   ·   " + modell["zusatz"] if modell.get("zusatz") else ""))
    y = Inches(1.95)
    for feat in modell["features"]:
        kachel(f, Inches(0.6), y, Inches(3.55), Inches(0.75), k["_f"]["flaeche"], rand=k["_f"]["linie"])
        kachel(f, Inches(0.6), y, Inches(0.07), Inches(0.75), k["_f"]["violett"], radius=0.0)
        textbox(f, Inches(0.85), y, Inches(3.2), Inches(0.75), feat["name"], 14,
                k["_f"]["violett_hell"], fett=True, anker=MSO_ANCHOR.MIDDLE, font=FWORD)
        textbox(f, Inches(4.55), y + Inches(0.08), Inches(8.2), Inches(0.7),
                feat["klartext"], 13.5, k["_f"]["text"])
        y += Inches(0.95)


def slide_rollen(prs, k):
    f = leere_folie(prs, k)
    r = k["rollen"]
    kopfzeile(f, k, r["titel"], r.get("untertitel"))
    b = Inches(3.9); x = Inches(0.6)
    for eintrag in r["punkte"]:
        kachel(f, x, Inches(2.1), b, Inches(3.4), k["_f"]["flaeche"], rand=k["_f"]["linie"])
        kachel(f, x, Inches(2.1), b, Inches(0.09), k["_f"]["violett"], radius=0.0)
        textbox(f, x + Inches(0.3), Inches(2.5), b - Inches(0.6), Inches(2.8), "", 14,
                k["_f"]["text"],
                zeilen=[(eintrag["wer"], 16, k["_f"]["violett_hell"], True, FWORD),
                        (eintrag["text"], 15, k["_f"]["text"], False, FBODY)])
        x += b + Inches(0.31)
    textbox(f, Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.6),
            r.get("fusszeile", ""), 13.5, k["_f"]["violett_hell"], fett=True)


def slide_abschluss(prs, k):
    f = leere_folie(prs, k)
    a = k["abschluss"]
    kachel(f, Inches(-1), Inches(1.3), Inches(15), Inches(3.0), None,
           grad=(k["_f"]["grund"], k["_f"]["violett_tief"]), gradwinkel=25, radius=0.0)
    kachel(f, 0, 0, BREITE, Inches(0.1), None, radius=0.0,
           grad=(k["_f"]["violett"], k["_f"]["violett_tief"]), gradwinkel=0)
    textbox(f, Inches(1.2), Inches(1.7), Inches(11), Inches(1.2),
            a["titel"], 42, k["_f"]["text"], font=FTITLE)
    kachel(f, Inches(1.25), Inches(3.05), Inches(1.4), Inches(0.05), k["_f"]["violett_hell"], radius=0.0)
    y = Inches(3.35)
    for punkt in a["punkte"]:
        kachel(f, Inches(1.2), y + Inches(0.06), Inches(0.1), Inches(0.5), k["_f"]["violett_hell"], radius=0.5)
        textbox(f, Inches(1.55), y, Inches(10.4), Inches(0.7), punkt, 17, k["_f"]["text"])
        y += Inches(0.85)
    textbox(f, Inches(1.2), Inches(6.4), Inches(10.8), Inches(0.6),
            a.get("fusszeile", ""), 13, k["_f"]["text_dim"], font=FLABEL)


def baue_deck(konfig_pfad, ziel_pfad):
    k = lade_konfig(konfig_pfad)
    prs = Presentation()
    prs.slide_width, prs.slide_height = BREITE, HOEHE
    slide_titel(prs, k)
    slide_momente(prs, k)
    slide_prinzipien(prs, k)
    if k.get("abteilungen"):
        slide_business(prs, k)
    if k.get("os_flow"):
        slide_os_flow(prs, k)
    for m in k.get("mockups", []):
        slide_mockup(prs, k, m)
    slide_schritte(prs, k)
    slide_modelle_visual(prs, k)
    for modell in k["modelle"]["liste"]:
        slide_klartext(prs, k, modell, k["modelle"]["preis_einheit"])
    slide_rollen(prs, k)
    slide_abschluss(prs, k)
    Path(ziel_pfad).parent.mkdir(parents=True, exist_ok=True)
    prs.save(ziel_pfad)
    print(f"OK: {ziel_pfad} ({len(prs.slides._sldIdLst)} Folien)")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        sys.exit("Aufruf: python baue_praesentation.py <konfig.json> <ziel.pptx>")
    baue_deck(sys.argv[1], sys.argv[2])
